from pptx import Presentation

from scripts.ingestion.models import AbstractIngestor, UnsupportedFileError


class PptxIngestor(AbstractIngestor):
    """
    Ingestor for PPTX files.
    """

    def ingest(self, filepath: str) -> list[tuple[str, dict]]:
        """
        Ingests data from the given PPTX filepath.

        Args:
            filepath: Path to the PPTX file to ingest.

        Returns:
            A list of tuples, where each tuple contains the extracted text
            and associated metadata (slide number, type, doc_type).
        """
        if not filepath.endswith(".pptx"):
            raise UnsupportedFileError("File is not a .pptx file.")

        extracted_data = []
        try:
            prs = Presentation(filepath)
            for i, slide in enumerate(prs.slides):
                slide_number = i + 1
                text_on_slide = []

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_on_slide.append(shape.text.strip())

                # Extract text from presenter notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        notes_meta = {
                            "slide_number": slide_number,
                            "type": "presenter_notes",
                            "doc_type": "pptx"
                        }
                        # Format notes text consistently
                        formatted_notes = (
                            f"Presenter Notes (Slide {slide_number}):\n"
                            f"{notes_text}"
                        )
                        extracted_data.append((formatted_notes, notes_meta))

                # Combine all text from the slide itself
                if text_on_slide:
                    # Filter out empty strings that might result from strip()
                    valid_texts_on_slide = [t for t in text_on_slide if t]
                    if valid_texts_on_slide:
                        slide_content = "\n".join(valid_texts_on_slide).strip()
                        # Ensure we don't add empty content after join and strip
                        if slide_content:
                            slide_meta = {
                                "slide_number": slide_number,
                                "type": "slide_content",
                                "doc_type": "pptx"
                            }
                            extracted_data.append((slide_content, slide_meta))

        except Exception as e:
            # Catch exceptions from python-pptx or other issues
            raise UnsupportedFileError(f"Error processing PPTX file {filepath}: {e}")

        return extracted_data
