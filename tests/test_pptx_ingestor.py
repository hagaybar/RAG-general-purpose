import unittest
import pathlib
# import os  # Unused import
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE

from scripts.ingestion.pptx import PptxIngestor
from scripts.ingestion.models import UnsupportedFileError


class TestPptxIngestor(unittest.TestCase):

    @classmethod
    def setUpClass(cls):
        cls.fixture_dir = pathlib.Path(__file__).parent / "fixtures" / "pptx"
        cls.fixture_dir.mkdir(parents=True, exist_ok=True)
        cls.sample_pptx_path = cls.fixture_dir / "test_presentation.pptx"
        cls.corrupted_pptx_path = cls.fixture_dir / "corrupted.pptx"
        cls.not_pptx_path = cls.fixture_dir / "not_a_pptx.txt"

        # Create a sample PPTX file if it doesn't exist
        if not cls.sample_pptx_path.exists():
            prs = Presentation()
            # Slide 1: Title, Textbox, Notes
            slide_layout_1 = prs.slide_layouts[0]  # Title slide layout
            slide_1 = prs.slides.add_slide(slide_layout_1)
            title_1 = slide_1.shapes.title
            title_1.text = "Slide 1 Title"

            left = top = width = height = Inches(1.0)
            tx_box_1 = slide_1.shapes.add_textbox(left, top + Inches(1.5),
                                                 width, height)
            tf_1 = tx_box_1.text_frame
            tf_1.text = "Slide 1 Textbox Content"

            notes_slide_1 = slide_1.notes_slide
            notes_tf_1 = notes_slide_1.notes_text_frame
            notes_tf_1.text = "Slide 1 Notes"

            # Slide 2: Content Placeholder, Shape, Notes
            # Blank layout (or one with content placeholder)
            slide_layout_2 = prs.slide_layouts[5]
            slide_2 = prs.slides.add_slide(slide_layout_2)

            # Add content placeholder text (if layout supports it,
            # otherwise textbox) For layout 5 (blank), we add a textbox.
            left_s2, top_s2 = Inches(1.0), Inches(1.0)
            width_s2, height_s2 = Inches(6.0), Inches(2.0)
            content_box_2 = slide_2.shapes.add_textbox(left_s2, top_s2,
                                                      width_s2, height_s2)
            content_tf_2 = content_box_2.text_frame
            content_tf_2.text = "Slide 2 Main Content"

            # Add a shape with text
            shape_2 = slide_2.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left_s2, top_s2 + Inches(2.5),
                Inches(3), Inches(1)
            )
            shape_tf_2 = shape_2.text_frame
            shape_tf_2.text = "Slide 2 Shape Text"

            notes_slide_2 = slide_2.notes_slide
            notes_tf_2 = notes_slide_2.notes_text_frame
            notes_tf_2.text = "Slide 2 Notes"

            # Slide 3: Only notes
            slide_layout_3 = prs.slide_layouts[5]  # Blank layout
            slide_3 = prs.slides.add_slide(slide_layout_3)
            notes_slide_3 = slide_3.notes_slide
            notes_tf_3 = notes_slide_3.notes_text_frame
            notes_tf_3.text = "Slide 3 Notes Only"

            prs.save(cls.sample_pptx_path)

        # Create a dummy corrupted pptx file
        with open(cls.corrupted_pptx_path, "w") as f:
            f.write("This is not a valid pptx content.")

        # Create a non-pptx file
        with open(cls.not_pptx_path, "w") as f:
            f.write("This is a text file.")


    def setUp(self):
        self.ingestor = PptxIngestor()

    def test_ingest_simple_presentation(self):
        extracted_data = self.ingestor.ingest(str(self.sample_pptx_path))

        self.assertIsNotNone(extracted_data)
        self.assertIsInstance(extracted_data, list)
        self.assertTrue(len(extracted_data) > 0, "Should extract some data")

        # PptxIngestor combines text from all shapes on a slide into a single entry.
        # The exact separator (\n or \n\n) might depend on how shape.text provides the text.
        # PptxIngestor now strips individual shape texts before joining with "\n".
        # This should result in a single "\n" between texts.
        expected_texts_and_types = {
            ("Slide 1 Title\nSlide 1 Textbox Content", "slide_content", 1),
            ("Presenter Notes (Slide 1):\nSlide 1 Notes",
             "presenter_notes", 1),
            ("Slide 2 Main Content\nSlide 2 Shape Text", "slide_content", 2),
            ("Presenter Notes (Slide 2):\nSlide 2 Notes",
             "presenter_notes", 2),
            ("Presenter Notes (Slide 3):\nSlide 3 Notes Only",
             "presenter_notes", 3),
        }

        extracted_details = set()

        # Debugging print can be uncommented if needed
        # print("\nExtracted Data for Debugging:")
        # for text_debug, meta_debug in extracted_data:
        #     print(f"  Text: {repr(text_debug)}, Meta: {meta_debug}")
        # print("---End Extracted Data---\n")

        for item in extracted_data:
            self.assertIsInstance(item, tuple)
            self.assertEqual(len(item), 2)
            text, meta = item
            self.assertIsInstance(text, str)
            self.assertIsInstance(meta, dict)

            self.assertEqual(meta.get('doc_type'), 'pptx')
            self.assertIn('slide_number', meta)
            self.assertIsInstance(meta['slide_number'], int)
            self.assertIn('type', meta)
            self.assertIn(meta['type'], ['slide_content', 'presenter_notes'])

            # Normalize extracted text for comparison if needed.
            normalized_text = text.strip()

            # Refined matching logic:
            # For presenter notes, expected text includes the prefix.
            # For slide content, it's the combined text from shapes.
            comparable_extracted_text = text  # Keep prefix for notes
            if meta['type'] == 'slide_content':
                # Use stripped text for slide content comparison
                comparable_extracted_text = normalized_text

            match_found_in_expected = False
            for expected_text, expected_type, expected_slide_num \
                    in expected_texts_and_types:
                if (comparable_extracted_text == expected_text and
                        meta['type'] == expected_type and
                        meta['slide_number'] == expected_slide_num):
                    extracted_details.add(
                        (expected_text, expected_type, expected_slide_num)
                    )
                    match_found_in_expected = True
                    break

            self.assertTrue(
                match_found_in_expected,
                f"Unexpected text segment: {repr(text)} with meta {meta}"
            )

        # Check if all expected items were found
        expected_count = len(expected_texts_and_types)
        actual_count = len(extracted_data)
        self.assertEqual(
            actual_count, expected_count,
            f"Expected {expected_count} segments, got {actual_count}.\n"
            f"Extracted: {[item[0] for item in extracted_data]}\n"
            f"Expected: {[item[0] for item in expected_texts_and_types]}"
        )
        self.assertEqual(
            extracted_details, expected_texts_and_types,
            "Not all expected text segments were found or matched correctly."
        )

    def test_ingest_non_pptx_file(self):
        with self.assertRaises(UnsupportedFileError):
            self.ingestor.ingest(str(self.not_pptx_path))

    def test_ingest_corrupted_pptx_file(self):
        # python-pptx might raise various errors for corrupted files,
        # PptxIngestor should wrap them in UnsupportedFileError.
        with self.assertRaises(UnsupportedFileError):
            self.ingestor.ingest(str(self.corrupted_pptx_path))

if __name__ == "__main__":
    unittest.main()
